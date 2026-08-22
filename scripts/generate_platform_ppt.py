#!/usr/bin/env python3
# ruff: noqa: E501, RUF001
"""Generate the RKNode platform introduction deck from the current implementation."""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "docs" / "ppt-assets"
OUTPUT = ROOT / "docs" / "RKNode平台介绍与节点数据流.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "Noto Sans CJK SC"
MONO = "Noto Sans Mono CJK SC"

WHITE = "FFFFFF"
INK = "18252B"
MUTED = "607078"
LIGHT = "F4F7F7"
LINE = "D8E1E1"
TEAL = "008A83"
TEAL_DARK = "006B66"
TEAL_LIGHT = "DDF3F0"
BLUE = "3777B7"
BLUE_LIGHT = "E6F0FA"
AMBER = "D58B24"
AMBER_LIGHT = "FFF0D9"
CORAL = "D95F4B"
CORAL_LIGHT = "FBE8E4"
GREEN = "4A8D63"
GREEN_LIGHT = "E7F3EA"
NAVY = "20383F"
DARK = "102126"
GRAY = "AAB7BA"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 16,
    color: str = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_VERTICAL_ANCHOR.TOP,
    font: str = FONT,
    margin: float = 0,
    fit: bool = False,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if fit else MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    p.line_spacing = 1.05
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return shape


def add_rich_text(slide, runs, x, y, w, h, *, align=PP_ALIGN.LEFT, valign=MSO_VERTICAL_ANCHOR.TOP):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    for text, size, color, bold in runs:
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return shape


def add_box(
    slide,
    x,
    y,
    w,
    h,
    *,
    fill=WHITE,
    line=LINE,
    radius=True,
    line_width=1,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(line_width)
    return shape


def add_circle(slide, x, y, d, *, fill=TEAL, line=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    return shape


def add_line(slide, x1, y1, x2, y2, *, color=LINE, width=1.5, dash=None):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = dash
    return line


def add_arrow(slide, x, y, w, h, *, fill=TEAL, direction="right"):
    shapes = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
    }
    arrow = slide.shapes.add_shape(shapes[direction], Inches(x), Inches(y), Inches(w), Inches(h))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = rgb(fill)
    arrow.line.color.rgb = rgb(fill)
    return arrow


def add_pill(slide, text, x, y, w, *, fill=TEAL_LIGHT, color=TEAL_DARK, size=10.5):
    box = add_box(slide, x, y, w, 0.34, fill=fill, line=fill, radius=True)
    add_text(
        slide,
        text,
        x,
        y,
        w,
        0.34,
        size=size,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_VERTICAL_ANCHOR.MIDDLE,
        fit=True,
    )
    return box


def add_bullets(slide, items, x, y, w, h, *, size=14, color=INK, accent=TEAL, gap=0.14):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_before = Pt(0)
        p.space_after = Pt(gap * 12)
        p.line_spacing = 1.08
        run = p.add_run()
        run.text = "● "
        run.font.name = FONT
        run.font.size = Pt(size - 2)
        run.font.color.rgb = rgb(accent)
        run = p.add_run()
        run.text = item
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(color)
    return shape


def add_step_number(slide, number, x, y, *, color=TEAL):
    add_circle(slide, x, y, 0.34, fill=color, line=color)
    add_text(
        slide,
        str(number),
        x,
        y,
        0.34,
        0.34,
        size=10,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_VERTICAL_ANCHOR.MIDDLE,
    )


def add_node_box(slide, title, subtitle, x, y, w, h, *, accent=TEAL, fill=WHITE, tag=None):
    add_box(slide, x, y, w, h, fill=fill, line=accent, line_width=1.5)
    add_box(slide, x, y, 0.09, h, fill=accent, line=accent, radius=False)
    add_text(slide, title, x + 0.25, y + 0.18, w - 0.45, 0.34, size=15, bold=True, color=INK)
    add_text(slide, subtitle, x + 0.25, y + 0.58, w - 0.45, h - 0.72, size=11.5, color=MUTED)
    if tag:
        add_pill(slide, tag, x + w - 1.18, y + 0.15, 0.92, fill=LIGHT, color=accent, size=9)


def add_icon_label(slide, code, title, subtitle, x, y, *, color=TEAL):
    add_circle(slide, x, y, 0.48, fill=color, line=color)
    add_text(
        slide,
        code,
        x,
        y,
        0.48,
        0.48,
        size=10,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_VERTICAL_ANCHOR.MIDDLE,
    )
    add_text(slide, title, x + 0.62, y - 0.01, 2.3, 0.28, size=13, bold=True)
    add_text(slide, subtitle, x + 0.62, y + 0.28, 2.3, 0.34, size=10.5, color=MUTED)


def add_picture_contain(slide, path: Path, x, y, w, h, *, border=True, bg=WHITE):
    add_box(slide, x, y, w, h, fill=bg, line=LINE if border else bg, radius=True)
    with Image.open(path) as image:
        iw, ih = image.size
    scale = min((w - 0.12) / iw, (h - 0.12) / ih)
    pw, ph = iw * scale, ih * scale
    px, py = x + (w - pw) / 2, y + (h - ph) / 2
    slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(pw), Inches(ph))


def add_slide_header(slide, title, index, *, kicker=None, dark=False):
    color = WHITE if dark else INK
    muted = "B7C8CB" if dark else MUTED
    add_box(slide, 0.48, 0.44, 0.08, 0.52, fill=TEAL, line=TEAL, radius=False)
    add_text(slide, title, 0.72, 0.34, 10.8, 0.55, size=25, color=color, bold=True)
    if kicker:
        add_text(slide, kicker, 0.74, 0.92, 10.6, 0.28, size=10.5, color=muted)
    add_text(slide, f"{index:02d}", 12.16, 0.41, 0.62, 0.32, size=10.5, color=muted, bold=True, align=PP_ALIGN.RIGHT)


def add_footer(slide, index, *, dark=False):
    color = "7F959A" if not dark else "85A0A5"
    add_line(slide, 0.55, 7.16, 12.78, 7.16, color="E2E8E8" if not dark else "294248", width=0.8)
    add_text(slide, "RKNode · RK3588 模型全生命周期平台", 0.58, 7.2, 4.4, 0.18, size=8.5, color=color)
    add_text(slide, f"{index} / 16", 11.95, 7.2, 0.8, 0.18, size=8.5, color=color, align=PP_ALIGN.RIGHT)


def new_slide(prs, *, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = rgb(bg)
    return slide


def add_flow_stage(slide, number, title, detail, x, y, w, *, color=TEAL, fill=WHITE):
    add_box(slide, x, y, w, 1.16, fill=fill, line=color, line_width=1.3)
    add_circle(slide, x + 0.18, y + 0.18, 0.36, fill=color, line=color)
    add_text(slide, str(number), x + 0.18, y + 0.18, 0.36, 0.36, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_VERTICAL_ANCHOR.MIDDLE)
    add_text(slide, title, x + 0.65, y + 0.15, w - 0.82, 0.32, size=14, bold=True)
    add_text(slide, detail, x + 0.2, y + 0.59, w - 0.4, 0.42, size=10.5, color=MUTED, fit=True)


def build_deck() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.title = "RKNode RK3588 模型全生命周期平台"
    prs.core_properties.subject = "中央控制面、直连节点与 RK3588 推理数据流"
    prs.core_properties.author = "RKNode Platform"
    prs.core_properties.keywords = (
        "RK3588, RKNN, direct node, Desired Revision, RKMPP, "
        "ByteTrack, Kafka, ZLM SEI"
    )

    # 1. Cover
    slide = new_slide(prs, bg=DARK)
    add_box(slide, 0, 0, 0.16, 7.5, fill=TEAL, line=TEAL, radius=False)
    for idx, width in enumerate((1.12, 0.9, 0.68)):
        add_box(slide, 0.84 + idx * 0.09, 0.8 + idx * 0.24, width, 0.16, fill=TEAL if idx == 0 else (BLUE if idx == 1 else AMBER), line=TEAL if idx == 0 else (BLUE if idx == 1 else AMBER), radius=False)
    add_text(slide, "RKNODE", 2.35, 0.78, 2.4, 0.4, size=15, color="AFC5C8", bold=True)
    add_text(slide, "RK3588 模型全生命周期平台", 0.86, 2.05, 11.1, 0.78, size=34, color=WHITE, bold=True)
    add_text(slide, "数据集 → 训练 → RKNN 转换 → 多板下发 → 板端业务推理", 0.9, 3.0, 10.8, 0.46, size=19, color="C4D4D6")
    add_line(slide, 0.9, 3.72, 6.85, 3.72, color=TEAL, width=2.5)
    labels = [
        ("TRAIN", "CPU / NVIDIA", BLUE),
        ("CONVERT", "RKNN Toolkit2", AMBER),
        ("DEPLOY", "Desired Revision", TEAL),
        ("INFER", "RK3588 C++ / NPU", CORAL),
    ]
    for idx, (top, bottom, color) in enumerate(labels):
        x = 0.9 + idx * 2.65
        add_box(slide, x, 4.35, 2.25, 1.02, fill=NAVY, line="355159", radius=True)
        add_text(slide, top, x + 0.2, 4.52, 1.85, 0.25, size=11, color=color, bold=True)
        add_text(slide, bottom, x + 0.2, 4.83, 1.85, 0.28, size=12, color=WHITE, bold=True, fit=True)
    add_text(slide, "重点：节点交互逻辑与数据流向", 0.9, 6.45, 6.5, 0.34, size=14, color="AFC5C8")
    add_text(slide, "基于当前项目实现 · 2026-08-20", 9.1, 6.45, 3.35, 0.34, size=11, color="7F989D", align=PP_ALIGN.RIGHT)

    # 2. Current product overview
    slide = new_slide(prs)
    add_slide_header(slide, "平台定位：中央控制面连接异构执行节点", 2, kicker="浏览器只访问中央平台；direct 节点通过独立服务接口接入")
    add_text(slide, "中央控制面", 0.72, 1.48, 2.3, 0.35, size=20, bold=True, color=TEAL_DARK)
    add_text(slide, "统一保存资产、任务、Release、期望状态和审计信息。", 0.72, 1.9, 3.55, 0.58, size=14, color=MUTED, fit=True)
    add_bullets(
        slide,
        [
            "Torch / Paddle 训练调度",
            "RKNN 转换与板端验证",
            "不可变 Release 与 Desired Revision",
            "多板收敛、业务推理与统一输出",
        ],
        0.72,
        2.64,
        3.55,
        2.35,
        size=13,
    )
    add_box(slide, 0.72, 5.42, 3.55, 0.84, fill=TEAL_LIGHT, line=TEAL_LIGHT)
    add_text(slide, "当前实现：训练、转换、推理三类节点均已直连接入", 0.92, 5.62, 3.15, 0.4, size=11.5, color=TEAL_DARK, bold=True, fit=True)
    add_picture_contain(slide, ASSETS / "overview.png", 4.55, 1.38, 8.16, 5.44)
    add_footer(slide, 2)

    # 3. Full lifecycle
    slide = new_slide(prs)
    add_slide_header(slide, "模型生命周期：从数据集到板端业务结果", 3, kicker="每个阶段都有明确输入、制品、验证和责任边界")
    stages = [
        ("01", "数据资产", "COCO / VOC / YOLO\n掩码 / PPOCR", BLUE),
        ("02", "模型训练", "PT / Pdparams\n+ ONNX + 日志", BLUE),
        ("03", "RKNN 转换", "RKNN + 验证报告\n+ 性能审计", AMBER),
        ("04", "模型发布", "不可变 Release\n标签 + 适配器契约", TEAL),
        ("05", "多板下发", "Desired Revision\n原子激活 + 回滚", TEAL),
        ("06", "业务推理", "RKMPP / RKNN\n分析 + 统一输出", CORAL),
    ]
    for idx, (num, title, detail, color) in enumerate(stages):
        x = 0.55 + idx * 2.1
        add_box(slide, x, 2.0, 1.72, 2.2, fill=WHITE, line=color, line_width=1.4)
        add_pill(slide, num, x + 0.18, 2.2, 0.54, fill=color, color=WHITE, size=9)
        add_text(slide, title, x + 0.18, 2.72, 1.36, 0.42, size=15, bold=True)
        add_text(slide, detail, x + 0.18, 3.28, 1.36, 0.64, size=10.5, color=MUTED, fit=True)
        if idx < len(stages) - 1:
            add_arrow(slide, x + 1.76, 2.85, 0.28, 0.38, fill=LINE)
    add_box(slide, 0.72, 4.85, 11.9, 1.24, fill=LIGHT, line=LIGHT)
    add_text(slide, "治理主线", 0.98, 5.07, 1.05, 0.28, size=12, bold=True, color=TEAL_DARK)
    add_text(slide, "数据可追溯  →  任务可观测  →  转换可验证  →  版本不可变  →  部署可回滚  →  结果可接入", 2.05, 5.02, 9.9, 0.38, size=15, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(slide, "训练与转换产物回到控制面；媒体由 ZLMediaKit 直达浏览器，不经过 Platform API。", 2.05, 5.52, 9.9, 0.3, size=11.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, 3)

    # 4. Deployment topology
    slide = new_slide(prs)
    add_slide_header(slide, "部署拓扑：控制面、媒体网关与四类执行节点", 4, kicker="Web/API 管控制；ZLMediaKit 管媒体；训练、转换和推理独立扩展")
    add_node_box(slide, "中心平台", "Web :5173 + Platform API\nZLMediaKit RTSP :8554 / WS-FLV :8081", 4.48, 1.35, 4.36, 1.26, accent=TEAL, fill=TEAL_LIGHT, tag="三镜像")
    add_node_box(slide, "训练节点 A", "CPU / NVIDIA CUDA\nTorch / Paddle 训练", 0.65, 3.85, 3.06, 1.37, accent=BLUE, tag=":10081")
    add_node_box(slide, "训练节点 N", "横向扩展 · 独立并发\n本地 /data/jobs", 0.65, 5.45, 3.06, 1.12, accent=BLUE, tag=":10081")
    add_node_box(slide, "RK3588 转换服务", "RKNN Toolkit2 2.3.2\nONNX → RKNN", 5.13, 3.85, 3.06, 1.37, accent=AMBER, tag=":10081")
    add_node_box(slide, "RK3588 推理服务", "C++ pipeline + RKNN Runtime\n原始编码 + schema-v2 SEI", 9.6, 3.85, 3.06, 1.37, accent=CORAL, tag=":10082")
    add_line(slide, 2.2, 3.7, 4.55, 2.67, color=BLUE, width=2)
    add_arrow(slide, 4.42, 2.47, 0.3, 0.38, fill=BLUE, direction="up")
    add_line(slide, 6.66, 3.68, 6.66, 2.68, color=AMBER, width=2)
    add_arrow(slide, 6.51, 2.47, 0.3, 0.38, fill=AMBER, direction="up")
    add_line(slide, 8.78, 2.67, 10.96, 3.7, color=CORAL, width=2)
    add_arrow(slide, 10.81, 3.51, 0.3, 0.38, fill=CORAL, direction="down")
    add_text(slide, "双向 API", 2.55, 3.23, 1.05, 0.25, size=9.5, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "双向 API", 6.02, 3.26, 1.28, 0.25, size=9.5, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "期望版本 / 状态", 9.0, 3.23, 1.42, 0.25, size=9.5, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 4.1, 5.75, 4.98, 0.66, fill=LIGHT, line=LIGHT)
    add_text(slide, "API 只做控制面；发布 RTSP 与播放 WS-FLV 使用不同短时 Token", 4.3, 5.92, 4.58, 0.28, size=11.5, color=INK, bold=True, align=PP_ALIGN.CENTER, fit=True)
    add_footer(slide, 4)

    # 5. Node onboarding and scheduling
    slide = new_slide(prs)
    add_slide_header(slide, "直连节点协议：身份、能力、健康、容量与鉴权", 5, kicker="direct 是新节点主线；旧 Worker / Agent 拉取模式仅保留兼容")
    columns = [
        ("1", "独立鉴权", "每节点独立 Token\n只授权匹配身份\n页面只写入不回显", BLUE),
        ("2", "健康身份", "GET /health\n协议 / name / kind\naccelerator / capabilities", TEAL),
        ("3", "容量匹配", "profile 与能力匹配\ncapacity = maxConcurrency\n− activeJobs", AMBER),
        ("4", "幂等执行", "dispatch jobId\nPUT desired revision\nDELETE cache", CORAL),
    ]
    for idx, (num, title, detail, color) in enumerate(columns):
        x = 0.68 + idx * 3.13
        add_box(slide, x, 1.68, 2.72, 2.48, fill=WHITE, line=color, line_width=1.4)
        add_circle(slide, x + 0.2, 1.9, 0.44, fill=color, line=color)
        add_text(slide, num, x + 0.2, 1.9, 0.44, 0.44, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_VERTICAL_ANCHOR.MIDDLE)
        add_text(slide, title, x + 0.76, 1.93, 1.7, 0.35, size=15, bold=True)
        add_text(slide, detail, x + 0.24, 2.55, 2.24, 1.25, size=11.5, color=MUTED, align=PP_ALIGN.CENTER, fit=True)
        if idx < 3:
            add_arrow(slide, x + 2.77, 2.68, 0.29, 0.4, fill=LINE)
    add_box(slide, 0.7, 4.67, 11.95, 1.34, fill=LIGHT, line=LIGHT)
    add_icon_label(slide, "T", "训练 / 转换任务", "中心仅推送 jobId；节点回源取数据", 1.0, 5.0, color=BLUE)
    add_icon_label(slide, "R", "推理部署", "中心推送期望版本；节点对齐实际版本", 4.68, 5.0, color=TEAL)
    add_icon_label(slide, "C", "缓存清理", "任务删除后下发 DELETE cache；离线后重试", 8.36, 5.0, color=CORAL)
    add_footer(slide, 5)

    # 6. Training matrix and direct dispatch
    slide = new_slide(prs)
    add_slide_header(slide, "训练调度：Torch / Paddle × CPU / CUDA", 6, kicker="统一任务契约覆盖两类框架和两类算力；节点无需共享中心磁盘")
    train_matrix = [
        ("Torch", "CPU", "YOLO / DeepLabV3+", "trainer-torch", BLUE),
        ("Torch", "CUDA", "YOLO / DeepLabV3+", "trainer-torch + CUDA overlay", BLUE),
        ("Paddle", "CPU", "PPOCR Det / Rec", "trainer-paddle", TEAL),
        ("Paddle", "CUDA", "PPOCR Det / Rec", "trainer-paddle + CUDA overlay", TEAL),
    ]
    for idx, (framework, accelerator, profiles, image, color) in enumerate(train_matrix):
        x = 0.68 + idx * 3.13
        add_box(slide, x, 1.42, 2.72, 2.05, fill=WHITE, line=color, line_width=1.35)
        add_pill(slide, accelerator, x + 0.2, 1.65, 0.82, fill=color, color=WHITE, size=9)
        add_text(slide, framework, x + 1.18, 1.62, 1.24, 0.34, size=17, bold=True, color=color, align=PP_ALIGN.RIGHT)
        add_text(slide, profiles, x + 0.22, 2.25, 2.28, 0.38, size=11.5, bold=True, align=PP_ALIGN.CENTER, fit=True)
        add_text(slide, image, x + 0.22, 2.82, 2.28, 0.34, size=9.5, color=MUTED, font=MONO, align=PP_ALIGN.CENTER, fit=True)
    add_text(slide, "direct 调度与数据回源", 0.72, 3.91, 2.55, 0.34, size=16, bold=True, color=TEAL_DARK)
    flow = [
        ("01", "创建任务", "数据集 + profile", TEAL),
        ("02", "下发 jobId", "direct dispatcher", BLUE),
        ("03", "授权回源", "数据集 / 初始权重", AMBER),
        ("04", "持续上报", "日志 / 指标 / 心跳", CORAL),
        ("05", "上传制品", "PT/Pdparams + ONNX", GREEN),
    ]
    for idx, (num, title, detail, color) in enumerate(flow):
        x = 0.68 + idx * 2.5
        add_box(slide, x, 4.42, 2.1, 1.35, fill=LIGHT, line=color, line_width=1.15)
        add_pill(slide, num, x + 0.17, 4.61, 0.5, fill=color, color=WHITE, size=8.5)
        add_text(slide, title, x + 0.78, 4.58, 1.1, 0.28, size=12.5, bold=True, fit=True)
        add_text(slide, detail, x + 0.18, 5.12, 1.74, 0.28, size=9.8, color=MUTED, align=PP_ALIGN.CENTER, fit=True)
        if idx < len(flow) - 1:
            add_arrow(slide, x + 2.14, 4.92, 0.27, 0.3, fill=LINE)
    add_box(slide, 0.72, 6.1, 11.86, 0.48, fill=BLUE_LIGHT, line=BLUE_LIGHT)
    add_text(slide, "节点 Token 只允许读取本节点当前任务引用的资产；任务临时区与容器日志均受生命周期和容量约束。", 0.92, 6.21, 11.46, 0.24, size=10.5, color=BLUE, bold=True, align=PP_ALIGN.CENTER, fit=True)
    add_footer(slide, 6)

    # 7. Dataset and training artifacts
    slide = new_slide(prs)
    add_slide_header(slide, "数据集与训练产物：格式、日志、指标、权重与 ONNX", 7, kicker="中心保存可追溯资产；训练节点只保留受控任务工作区")
    add_box(slide, 0.68, 1.45, 3.2, 4.92, fill=BLUE_LIGHT, line=BLUE)
    add_text(slide, "数据集输入", 0.96, 1.73, 2.64, 0.36, size=17, bold=True, color=BLUE)
    add_text(slide, "上传后完成格式探测、清单生成、类别统计与可训练性检查", 0.96, 2.18, 2.64, 0.64, size=11.2, color=MUTED, fit=True)
    for idx, label in enumerate(("COCO", "VOC", "YOLO", "MASK", "PPOCR")):
        row, col = divmod(idx, 2)
        add_pill(slide, label, 0.96 + col * 1.34, 3.05 + row * 0.55, 1.12, fill=WHITE, color=BLUE, size=9.5)
    add_box(slide, 0.96, 4.92, 2.64, 0.92, fill=WHITE, line=WHITE)
    add_text(slide, "manifest", 1.18, 5.1, 0.9, 0.24, size=11, color=BLUE, bold=True, font=MONO)
    add_text(slide, "类别 / 划分 / SHA", 1.18, 5.42, 2.0, 0.24, size=10, color=MUTED)

    add_box(slide, 4.28, 1.45, 3.2, 4.92, fill=LIGHT, line=TEAL)
    add_text(slide, "训练执行", 4.56, 1.73, 2.64, 0.36, size=17, bold=True, color=TEAL_DARK)
    train_steps = [
        ("01", "解析任务快照", "数据与超参数固定"),
        ("02", "框架训练", "Torch / Paddle"),
        ("03", "实时遥测", "日志 / 进度 / 指标"),
        ("04", "静态导出", "ONNX + manifest"),
    ]
    for idx, (num, title, detail) in enumerate(train_steps):
        y = 2.35 + idx * 0.83
        add_pill(slide, num, 4.58, y, 0.54, fill=TEAL, color=WHITE, size=8.5)
        add_text(slide, title, 5.3, y - 0.02, 1.5, 0.25, size=11.5, bold=True)
        add_text(slide, detail, 5.3, y + 0.3, 1.64, 0.22, size=9.5, color=MUTED, fit=True)
    add_box(slide, 4.56, 5.75, 2.64, 0.38, fill=TEAL_LIGHT, line=TEAL_LIGHT)
    add_text(slide, "删除 / 重试均保留审计链路", 4.7, 5.83, 2.36, 0.2, size=9.5, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER)

    add_box(slide, 7.88, 1.45, 4.77, 4.92, fill=WHITE, line=LINE)
    add_text(slide, "可追溯产物", 8.18, 1.73, 4.15, 0.36, size=17, bold=True, color=INK)
    artifacts = [
        ("权重", "PT / Pdparams", BLUE),
        ("部署中间件", "静态 ONNX", TEAL),
        ("过程证据", "日志 / 指标 / 事件", AMBER),
        ("模型契约", "输入输出 / 标签 / 预处理", CORAL),
    ]
    for idx, (title, detail, color) in enumerate(artifacts):
        y = 2.35 + idx * 0.82
        add_box(slide, 8.18, y, 4.15, 0.64, fill=LIGHT, line=color, line_width=1.05)
        add_pill(slide, title, 8.38, y + 0.15, 1.08, fill=color, color=WHITE, size=8.7)
        add_text(slide, detail, 9.68, y + 0.13, 2.36, 0.28, size=10.8, color=INK, bold=True, valign=MSO_VERTICAL_ANCHOR.MIDDLE, fit=True)
    add_box(slide, 8.18, 5.82, 4.15, 0.31, fill=GREEN_LIGHT, line=GREEN_LIGHT)
    add_text(slide, "中心长期保存 · 节点缓存按任务回收", 8.38, 5.87, 3.75, 0.2, size=9.4, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 3.92, 3.55, 0.28, 0.38, fill=LINE)
    add_arrow(slide, 7.52, 3.55, 0.28, 0.38, fill=LINE)
    add_footer(slide, 7)

    # 8. Conversion sequence
    slide = new_slide(prs)
    add_slide_header(slide, "RKNN 转换：Toolkit2、量化、板端验证与部署就绪", 8, kicker="转换节点在 RK3588 容器内使用 Toolkit2 / Runtime 完成真实板端资格判定")
    steps = [
        (1, "选择输入", "训练 ONNX\n+ 输入尺度 / 量化策略", BLUE),
        (2, "回源下载", "ONNX + 校准数据\n+ 标签 / manifest", TEAL),
        (3, "构建", "Toolkit2 build\nexport RKNN", AMBER),
        (4, "板端验证", "init_runtime(rk3588)\n固定样本推理", CORAL),
        (5, "性能审计", "warmup 5 / timing 20\nCPU fallback 检查", CORAL),
        (6, "回传产物", "RKNN + 日志\n报告 + readiness", TEAL),
    ]
    for idx, (num, title, detail, color) in enumerate(steps):
        x = 0.52 + idx * 2.13
        add_flow_stage(slide, num, title, detail, x, 1.65, 1.73, color=color)
        if idx < 5:
            add_arrow(slide, x + 1.77, 2.05, 0.28, 0.35, fill=LINE)
    add_box(slide, 0.72, 3.55, 5.7, 2.18, fill=AMBER_LIGHT, line=AMBER_LIGHT)
    add_text(slide, "精度优先的转换约束", 1.0, 3.84, 2.55, 0.35, size=16, bold=True, color="9B5F12")
    add_bullets(
        slide,
        ["按模型类型选择官方兼容配置", "INT8 使用代表性校准集", "真实 runtime 验证通过后才可发布"],
        1.0,
        4.32,
        5.0,
        1.16,
        size=12,
        accent=AMBER,
    )
    add_box(slide, 6.76, 3.55, 5.83, 2.18, fill=LIGHT, line=LIGHT)
    add_text(slide, "转换成功不等于部署就绪", 7.06, 3.84, 3.2, 0.35, size=16, bold=True, color=TEAL_DARK)
    add_bullets(
        slide,
        ["deploymentReady / performanceReady", "输入输出张量契约与后处理适配器", "校验统计、耗时、CPU 回退和失败原因"],
        7.06,
        4.32,
        5.0,
        1.16,
        size=12,
    )
    add_footer(slide, 8)

    # 9. Release boundary
    slide = new_slide(prs)
    add_slide_header(slide, "不可变 Release：制品、校验和与运行契约", 9, kicker="推理任务绑定已发布版本，不直接引用任意 RKNN 文件")
    add_node_box(slide, "训练 / 转换溯源", "dataset · train job\nconversion job · source ONNX", 0.72, 2.0, 3.05, 1.42, accent=BLUE, fill=BLUE_LIGHT)
    add_arrow(slide, 3.94, 2.48, 0.64, 0.4, fill=LINE)
    add_box(slide, 4.72, 1.55, 4.06, 2.34, fill=TEAL_LIGHT, line=TEAL, line_width=1.6)
    add_pill(slide, "IMMUTABLE", 5.02, 1.82, 1.2, fill=TEAL, color=WHITE, size=9)
    add_text(slide, "Model Release", 5.02, 2.28, 3.05, 0.45, size=22, bold=True, color=TEAL_DARK)
    add_text(slide, "主模型 + 全部二级模型 + SHA256 + manifest\n输入输出契约 + 验证报告 + 适配器 + 运行参数", 5.02, 2.86, 3.2, 0.7, size=11.5, color=MUTED, fit=True)
    add_arrow(slide, 8.92, 2.48, 0.64, 0.4, fill=LINE)
    add_node_box(slide, "部署与回滚", "任务绑定精确 releaseId\n版本历史可审计、可恢复", 9.7, 2.0, 2.9, 1.42, accent=CORAL, fill=CORAL_LIGHT)
    add_text(slide, "后处理适配器", 0.74, 4.56, 2.0, 0.35, size=15, bold=True)
    adapters = [
        ("YOLO", "DFL / Anchored / v6 / v7 / v10", BLUE),
        ("DeepLab", "semantic logits → mask", TEAL),
        ("PPOCR", "DB detection / CTC recognition", AMBER),
    ]
    for idx, (name, desc, color) in enumerate(adapters):
        x = 0.74 + idx * 4.02
        add_box(slide, x, 5.05, 3.68, 1.04, fill=WHITE, line=color, line_width=1.2)
        add_pill(slide, name, x + 0.18, 5.23, 0.88, fill=color, color=WHITE, size=9)
        add_text(slide, desc, x + 1.18, 5.18, 2.2, 0.56, size=10.5, color=MUTED, bold=True, valign=MSO_VERTICAL_ANCHOR.MIDDLE, fit=True)
    add_footer(slide, 9)

    # 10. Multi-board deployment
    slide = new_slide(prs)
    add_slide_header(slide, "Desired Revision：多板下发、原子激活与回滚", 10, kicker="中心维护 desired，节点上报 actual；离线恢复后继续收敛")
    add_node_box(slide, "已发布模型", "Release v12\nSHA256 / manifest / adapter", 0.62, 1.45, 2.55, 1.25, accent=TEAL, fill=TEAL_LIGHT)
    add_node_box(slide, "部署批次", "canary / rolling / all_at_once\n目标板卡 + NPU 核策略", 3.53, 1.45, 3.12, 1.25, accent=BLUE, fill=BLUE_LIGHT)
    add_node_box(slide, "期望版本", "desiredRevision = 28\n每个节点独立目标状态", 7.02, 1.45, 2.68, 1.25, accent=AMBER, fill=AMBER_LIGHT)
    add_arrow(slide, 3.18, 1.87, 0.28, 0.38, fill=LINE)
    add_arrow(slide, 6.67, 1.87, 0.28, 0.38, fill=LINE)
    boards = [
        ("板卡 A", "actual = 28", "健康", GREEN),
        ("板卡 B", "actual = 27 → 28", "更新中", AMBER),
        ("板卡 C", "offline · pending", "待重连", CORAL),
    ]
    for idx, (name, rev, state, color) in enumerate(boards):
        y = 3.28 + idx * 1.02
        add_box(slide, 8.9, y, 3.4, 0.76, fill=WHITE, line=color, line_width=1.3)
        add_text(slide, name, 9.16, y + 0.16, 0.8, 0.25, size=12, bold=True)
        add_text(slide, rev, 10.0, y + 0.16, 1.33, 0.25, size=10.5, color=MUTED, font=MONO, fit=True)
        add_pill(slide, state, 11.4, y + 0.19, 0.64, fill=color, color=WHITE, size=8.5)
        add_arrow(slide, 7.46, y + 0.21, 1.05, 0.32, fill=color)
    add_box(slide, 0.72, 3.18, 6.15, 2.96, fill=LIGHT, line=LIGHT)
    add_text(slide, "节点执行阶段", 1.0, 3.48, 1.75, 0.34, size=16, bold=True, color=TEAL_DARK)
    stages = ["下载全部制品", "SHA256 校验", "静态契约检查", "Runtime 探测", "模型预热", "排空旧流", "原子激活", "健康上报"]
    for idx, label in enumerate(stages):
        row, col = divmod(idx, 4)
        x = 1.0 + col * 1.38
        y = 4.08 + row * 0.75
        add_pill(slide, label, x, y, 1.18, fill=TEAL_LIGHT if idx < 6 else GREEN_LIGHT, color=TEAL_DARK if idx < 6 else GREEN, size=8.8)
        if col < 3:
            add_arrow(slide, x + 1.22, y + 0.07, 0.12, 0.2, fill=LINE)
    add_text(slide, "候选目录失败即丢弃；current / previous 保持完整", 1.0, 5.7, 5.4, 0.28, size=11, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 10)

    # 11. RK3588 data plane
    slide = new_slide(prs)
    add_slide_header(slide, "RK3588 数据面：RKMPP、RKNN 与 ByteTrack", 11, kicker="Node Service 协调版本；C++ pipeline 承担硬解码、NPU 推理和实时业务处理")
    add_box(slide, 0.66, 1.36, 12.02, 1.62, fill=LIGHT, line=LINE)
    add_pill(slide, "CONTROL", 0.92, 1.62, 0.96, fill=TEAL, color=WHITE, size=8.8)
    control_nodes = [
        ("Node Service", "desired / actual\n下载、校验、激活、上报", TEAL),
        ("模型仓库", "staging / current / previous\n内容寻址缓存", AMBER),
        ("Runtime Adapter", "生成配置 · 进程生命周期\n健康探针与失败恢复", BLUE),
    ]
    for idx, (title, detail, color) in enumerate(control_nodes):
        x = 2.12 + idx * 3.38
        add_node_box(slide, title, detail, x, 1.58, 2.92, 1.06, accent=color, fill=WHITE)
        if idx < len(control_nodes) - 1:
            add_arrow(slide, x + 2.98, 1.95, 0.28, 0.32, fill=LINE)

    add_pill(slide, "C++ DATA PLANE", 0.72, 3.35, 1.55, fill=CORAL, color=WHITE, size=8.8)
    pipeline = [
        ("RKMPP / OpenCV", "RTSP 硬解码\n或旧任务兼容", BLUE),
        ("Primary RKNN", "YOLO / Seg / OCR\n主模型推理", CORAL),
        ("ByteTrack", "稳定 track ID\n检测任务可选", AMBER),
        ("Secondary RKNN", "目标裁剪推理\n父目标关联", CORAL),
        ("Analytics", "区域 / 越线\n事件状态", TEAL),
        ("SEI / Event", "原码流元数据\n事件状态", AMBER),
        ("Outputs", "schema v2\n多路输出", GREEN),
    ]
    for idx, (title, detail, color) in enumerate(pipeline):
        x = 0.52 + idx * 1.79
        add_box(slide, x, 3.92, 1.48, 1.55, fill=WHITE, line=color, line_width=1.2)
        add_text(slide, title, x + 0.12, 4.15, 1.24, 0.34, size=10.5, bold=True, color=INK, align=PP_ALIGN.CENTER, fit=True)
        add_text(slide, detail, x + 0.12, 4.7, 1.24, 0.48, size=8.8, color=MUTED, align=PP_ALIGN.CENTER, fit=True)
        if idx < len(pipeline) - 1:
            add_arrow(slide, x + 1.51, 4.53, 0.22, 0.28, fill=color)
    add_box(slide, 0.72, 5.84, 12.0, 0.68, fill=CORAL_LIGHT, line=CORAL_LIGHT)
    add_text(slide, "RKMPP：RTSP H.264/H.265 · NPU Core 0/1/2：shared 或 exclusive · 未启用的业务节点从运行图中省略", 0.96, 6.03, 11.52, 0.28, size=10.8, color=CORAL, bold=True, align=PP_ALIGN.CENTER, fit=True)
    add_footer(slide, 11)

    # 12. Business analytics
    slide = new_slide(prs)
    add_slide_header(slide, "业务分析：二级推理、区域、越线与浏览器叠加", 12, kicker="结构化结果进入 SEI；Canvas 按视频 PTS 绘制，不重编码视频")
    add_picture_contain(slide, ASSETS / "inference.png", 0.62, 1.42, 5.26, 3.48)
    add_box(slide, 0.72, 5.13, 5.06, 0.92, fill=GREEN_LIGHT, line=GREEN_LIGHT)
    add_text(slide, "当前节点已上报", 0.96, 5.33, 1.48, 0.25, size=11, color=GREEN, bold=True)
    add_text(slide, "RKMPP · ByteTrack · Analytics · Event · Secondary", 2.38, 5.28, 3.14, 0.36, size=9.7, color=GREEN, bold=True, align=PP_ALIGN.RIGHT, fit=True)

    business = [
        ("二级推理", "普通非人脸 YOLO；裁剪主目标后推理，使用 parentTrackId / parentDetectionIndex 关联。", CORAL),
        ("区域计数", "检测框底边中心点进入多边形；按 track ID 维护进入/离开与 holdFrames 防抖。", TEAL),
        ("越线检测", "连续轨迹跨越有向线段；按 track ID 与方向去重并累计双向计数。", BLUE),
        ("Canvas 叠加", "目标框、分割掩码、OCR、轨迹和计数由浏览器按同一视频 PTS 绘制。", AMBER),
    ]
    for idx, (title, detail, color) in enumerate(business):
        y = 1.44 + idx * 1.13
        add_box(slide, 6.2, y, 6.46, 0.94, fill=WHITE, line=color, line_width=1.15)
        add_pill(slide, title, 6.46, y + 0.27, 1.28, fill=color, color=WHITE, size=9)
        add_text(slide, detail, 7.98, y + 0.16, 4.36, 0.58, size=10.1, color=MUTED, valign=MSO_VERTICAL_ANCHOR.MIDDLE, fit=True)
    add_box(slide, 6.2, 6.0, 6.46, 0.5, fill=LIGHT, line=LIGHT)
    add_text(slide, "track ID → 区域/越线状态 → schema-v2 SEI → Canvas", 6.46, 6.12, 5.94, 0.25, size=10.4, color=INK, bold=True, align=PP_ALIGN.CENTER, fit=True)
    add_footer(slide, 12)

    # 13. Unified outputs
    slide = new_slide(prs)
    add_slide_header(slide, "统一输出：业务 Sink 与 RTSP + SEI 浏览器直连", 13, kicker="API 不传媒体字节；ZLMediaKit 无重编码 remux 为 WS-FLV")
    add_node_box(slide, "FrameResult", "检测 / 分割 / OCR / 跟踪\n区域 / 越线 / 事件 / 媒体", 0.72, 1.48, 3.1, 1.25, accent=CORAL, fill=CORAL_LIGHT)
    add_arrow(slide, 3.96, 1.91, 0.56, 0.38, fill=LINE)
    add_box(slide, 4.65, 1.35, 4.06, 1.5, fill=TEAL_LIGHT, line=TEAL, line_width=1.5)
    add_pill(slide, "SCHEMA V2", 4.94, 1.62, 1.1, fill=TEAL, color=WHITE, size=8.8)
    add_text(slide, "统一结构化协议", 4.94, 2.06, 3.4, 0.34, size=18, bold=True, color=TEAL_DARK)
    add_text(slide, "task / revision / timestamp / sequence + results / analytics / media", 4.94, 2.47, 3.4, 0.22, size=8.8, color=MUTED, font=MONO, fit=True)
    add_arrow(slide, 8.84, 1.91, 0.56, 0.38, fill=LINE)
    add_node_box(slide, "统一序列化器", "保留旧字段兼容\n所有结果通道字段一致", 9.54, 1.48, 3.06, 1.25, accent=BLUE, fill=BLUE_LIGHT)

    outputs = [
        ("JSONL", "板端持久化", "任务目录", TEAL),
        ("HTTP", "业务系统推送", "结构化 API", BLUE),
        ("Kafka", "异步有界消息", "独立队列", AMBER),
        ("ZLM SEI", "原码流注入", "无视频重编码", CORAL),
        ("Browser", "WS-FLV + Canvas", "直接访问网关", GREEN),
    ]
    add_line(slide, 6.67, 2.92, 6.67, 3.43, color=LINE, width=1.5)
    add_line(slide, 1.7, 3.43, 11.7, 3.43, color=LINE, width=1.5)
    for idx, (title, purpose, owner, color) in enumerate(outputs):
        x = 0.58 + idx * 2.52
        add_line(slide, x + 1.0, 3.43, x + 1.0, 3.72, color=LINE, width=1.5)
        add_box(slide, x, 3.72, 2.02, 1.62, fill=WHITE, line=color, line_width=1.25)
        add_pill(slide, title, x + 0.24, 3.98, 1.54, fill=color, color=WHITE, size=9)
        add_text(slide, purpose, x + 0.18, 4.54, 1.66, 0.28, size=10.6, bold=True, align=PP_ALIGN.CENTER, fit=True)
        add_text(slide, owner, x + 0.18, 4.92, 1.66, 0.22, size=9.2, color=MUTED, align=PP_ALIGN.CENTER, fit=True)
    add_box(slide, 0.72, 5.78, 12.0, 0.72, fill=LIGHT, line=LIGHT)
    add_text(slide, "RK3588 → 鉴权 RTSP :8554 → ZLMediaKit → 鉴权 WS-FLV :8081 → video + Canvas；Platform API 仅签发凭据", 0.96, 5.98, 11.52, 0.3, size=10.8, color=INK, bold=True, align=PP_ALIGN.CENTER, fit=True)
    add_footer(slide, 13)

    # 14. Failure recovery
    slide = new_slide(prs)
    add_slide_header(slide, "可靠性：幂等、缓存清理、恢复与最终收敛", 14, kicker="允许节点离线；调度、激活、删除和重连都具备明确恢复语义")
    cases = [
        ("节点离线", "任务保持排队或部署保持 pending；恢复后重新探测并继续收敛。", BLUE, "RETRY"),
        ("任务失败", "保留日志、指标和错误详情；前端提供重新训练 / 重新转换。", AMBER, "RERUN"),
        ("Revision 失败", "候选版本不切换；节点保留 current / previous 并恢复上一健康进程。", CORAL, "ROLLBACK"),
        ("任务删除", "中心删除记录并创建 node_cleanup；在线立即清理，离线后重试。", TEAL, "CLEANUP"),
    ]
    for idx, (title, body, color, code) in enumerate(cases):
        row, col = divmod(idx, 2)
        x = 0.72 + col * 6.15
        y = 1.5 + row * 2.15
        add_box(slide, x, y, 5.72, 1.75, fill=WHITE, line=color, line_width=1.4)
        add_pill(slide, code, x + 0.25, y + 0.25, 1.04, fill=color, color=WHITE, size=8.8)
        add_text(slide, title, x + 1.5, y + 0.22, 3.7, 0.35, size=16, bold=True)
        add_text(slide, body, x + 0.25, y + 0.84, 5.15, 0.62, size=11.5, color=MUTED, fit=True)
    add_box(slide, 0.72, 5.96, 11.86, 0.56, fill=LIGHT, line=LIGHT)
    add_text(slide, "health identity 校验 · capacity 计算 · 幂等 dispatch · desired/actual 对账 · node_cleanup 重试 · 有界容器日志", 0.92, 6.1, 11.46, 0.27, size=10.8, color=INK, bold=True, align=PP_ALIGN.CENTER, fit=True)
    add_footer(slide, 14)

    # 15. Deployment and security matrix
    slide = new_slide(prs)
    add_slide_header(slide, "部署与安全：在线/离线镜像、Compose、Token 与隧道", 15, kicker="同一版本可在联网与隔离网络交付；跨网段仍保持中心主动访问节点")
    add_text(slide, "部署矩阵", 0.72, 1.42, 1.4, 0.32, size=16, bold=True, color=TEAL_DARK)
    deployments = [
        ("在线镜像", "按版本构建 / 拉取", "联网环境", TEAL),
        ("完全离线", "镜像 tar + Compose + 校验清单", "隔离环境", AMBER),
        ("中央平台", "API / Web / Media · 两个离线包", "5173 / 8554 / 8081", TEAL),
        ("训练节点", "CPU 基础 + CUDA overlay", "Torch / Paddle · 10081", BLUE),
        ("RK3588", "转换 / 推理独立容器", "10081 / 10082", CORAL),
    ]
    for idx, (name, detail, scope, color) in enumerate(deployments):
        y = 1.87 + idx * 0.88
        add_box(slide, 0.72, y, 6.0, 0.7, fill=LIGHT, line=color, line_width=1.05)
        add_pill(slide, name, 0.94, y + 0.18, 1.12, fill=color, color=WHITE, size=8.7)
        add_text(slide, detail, 2.28, y + 0.12, 2.68, 0.24, size=10.3, bold=True, fit=True)
        add_text(slide, scope, 4.98, y + 0.12, 1.48, 0.34, size=9.2, color=MUTED, align=PP_ALIGN.RIGHT, fit=True)

    add_text(slide, "鉴权与跨网", 7.08, 1.42, 1.6, 0.32, size=16, bold=True, color=TEAL_DARK)
    security = [
        ("三类独立 Token", "节点长期 / 发布短时 / 播放 60 秒", TEAL),
        ("防火墙收口", "仅允许中央访问节点服务端口", BLUE),
        ("VPN", "推荐长期跨网段 direct 访问", GREEN),
        ("SSH 隧道", "适合临时交付与受限网络", AMBER),
        ("显式设备权限", "NPU / MPP / RGA / dma-heap；不使用 privileged", CORAL),
    ]
    for idx, (title, detail, color) in enumerate(security):
        y = 1.87 + idx * 0.88
        add_circle(slide, 7.12, y + 0.16, 0.34, fill=color, line=color)
        add_text(slide, str(idx + 1), 7.12, y + 0.16, 0.34, 0.34, size=8.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_VERTICAL_ANCHOR.MIDDLE)
        add_text(slide, title, 7.65, y + 0.08, 1.72, 0.28, size=11.2, bold=True, fit=True)
        add_text(slide, detail, 9.5, y + 0.08, 2.86, 0.42, size=9.8, color=MUTED, fit=True)
    add_box(slide, 0.72, 6.38, 11.94, 0.3, fill=CORAL_LIGHT, line=CORAL_LIGHT)
    add_text(slide, "转换与推理可同板部署，但容器、端口、数据卷、Token 和日志边界相互独立", 0.96, 6.42, 11.46, 0.2, size=9.8, color=CORAL, bold=True, align=PP_ALIGN.CENTER, fit=True)
    add_footer(slide, 15)

    # 16. Value and boundary
    slide = new_slide(prs, bg=DARK)
    add_slide_header(slide, "核心价值与技术边界", 16, kicker="一套控制面连接异构训练算力与多台 RK3588 执行节点", dark=True)
    values = [
        ("统一治理", "数据、训练、转换、Release、\n部署和推理结果形成可追溯链路", TEAL),
        ("异构调度", "direct 连接 Torch/Paddle、\nCPU/CUDA 和 RK3588", BLUE),
        ("版本可信", "SHA256 + Desired Revision\n候选全验、原子激活、失败回滚", AMBER),
        ("业务闭环", "RKMPP、RKNN、ByteTrack\n业务分析与五路统一输出", CORAL),
    ]
    for idx, (title, body, color) in enumerate(values):
        x = 0.7 + idx * 3.13
        add_box(slide, x, 1.62, 2.72, 2.86, fill=NAVY, line="355159")
        add_circle(slide, x + 0.24, 1.92, 0.48, fill=color, line=color)
        add_text(slide, f"0{idx + 1}", x + 0.24, 1.92, 0.48, 0.48, size=9.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_VERTICAL_ANCHOR.MIDDLE)
        add_text(slide, title, x + 0.24, 2.62, 2.18, 0.45, size=19, color=WHITE, bold=True)
        add_text(slide, body, x + 0.24, 3.26, 2.18, 0.82, size=10.7, color="B7C8CB", fit=True)
    add_box(slide, 0.7, 4.92, 12.02, 1.34, fill="162E34", line="3A5359")
    add_pill(slide, "ACCEPTANCE", 0.98, 5.18, 1.28, fill=CORAL, color=WHITE, size=8.8)
    add_text(slide, "目标环境最终验收", 2.5, 5.15, 2.2, 0.3, size=14.5, color=WHITE, bold=True)
    add_text(slide, "至少一条真实 RTSP 完成 RKMPP、NPU 推理、原始 H.264 + SEI 发布、WS-FLV 播放和 Canvas 叠加；外部 Kafka 与 ZLMediaKit 端到端联调。", 2.5, 5.55, 9.72, 0.46, size=10.5, color="B7C8CB", fit=True)
    add_text(slide, "RKNode Platform", 0.76, 6.78, 2.2, 0.25, size=10.5, color="789297", bold=True)
    add_text(slide, "END", 11.76, 6.78, 0.7, 0.25, size=10.5, color="789297", bold=True, align=PP_ALIGN.RIGHT)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build_deck()
    print(path)
