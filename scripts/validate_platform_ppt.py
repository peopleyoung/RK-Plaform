#!/usr/bin/env python3
# ruff: noqa: RUF001
"""Validate the generated RKNode technical deck against its content contract."""

from __future__ import annotations

import argparse
import re
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_DECK = ROOT / "docs" / "RKNode平台介绍与节点数据流.pptx"

EXPECTED_TITLES = [
    "RK3588 模型全生命周期平台",
    "平台定位：中央控制面连接异构执行节点",
    "模型生命周期：从数据集到板端业务结果",
    "部署拓扑：控制面、媒体网关与四类执行节点",
    "直连节点协议：身份、能力、健康、容量与鉴权",
    "训练调度：Torch / Paddle × CPU / CUDA",
    "数据集与训练产物：格式、日志、指标、权重与 ONNX",
    "RKNN 转换：Toolkit2、量化、板端验证与部署就绪",
    "不可变 Release：制品、校验和与运行契约",
    "Desired Revision：多板下发、原子激活与回滚",
    "RK3588 数据面：RKMPP、RKNN 与 ByteTrack",
    "业务分析：二级推理、区域、越线与浏览器叠加",
    "统一输出：业务 Sink 与 RTSP + SEI 浏览器直连",
    "可靠性：幂等、缓存清理、恢复与最终收敛",
    "部署与安全：在线/离线镜像、Compose、Token 与隧道",
    "核心价值与技术边界",
]

REQUIRED_TERMS = [
    "direct",
    "Torch",
    "Paddle",
    "CUDA",
    "RKNN Toolkit2",
    "SHA256",
    "Desired Revision",
    "RKMPP",
    "ByteTrack",
    "二级推理",
    "区域",
    "越线",
    "Kafka",
    "ZLM SEI",
    "WS-FLV",
    "5173",
    "8554",
    "8081",
    "两个离线包",
    "发布短时",
    "播放 60 秒",
    "完全离线",
    "VPN",
    "SSH 隧道",
    "目标环境最终验收",
]

FORBIDDEN_TERMS = [
    "2026-08-12",
    "2026-08-14",
    "中心平台固定端口 8080",
    "浏览器直连板卡",
    "privileged: true",
    "人脸识别",
    "dev-admin-token",
    "限频 JPEG",
    "中心 latest",
    "短时预览会话",
]

IPV4_PATTERN = re.compile(
    r"(?<![\d.])(?:25[0-5]|2[0-4]\d|1?\d?\d)"
    r"(?:\.(?:25[0-5]|2[0-4]\d|1?\d?\d)){3}(?![\d.])"
)
SECRET_PATTERNS = [
    ("Bearer token", re.compile(r"\bBearer\s+\S+", re.IGNORECASE)),
    ("private key", re.compile(r"BEGIN [A-Z ]*PRIVATE KEY")),
]


def slide_texts(prs: Presentation) -> list[str]:
    """Return normalized text for every slide in display order."""
    texts: list[str] = []
    for slide in prs.slides:
        parts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text")]
        texts.append("\n".join(part for part in parts if part))
    return texts


def validate_deck(path: Path) -> list[str]:
    """Return human-readable contract violations for one PPTX file."""
    prs = Presentation(path)
    texts = slide_texts(prs)
    problems: list[str] = []

    if len(prs.slides) != len(EXPECTED_TITLES):
        problems.append(f"slide count: expected {len(EXPECTED_TITLES)}, got {len(prs.slides)}")

    for index, expected in enumerate(EXPECTED_TITLES):
        if index >= len(texts):
            break
        if expected not in texts[index]:
            problems.append(f"slide {index + 1}: missing title {expected!r}")

    all_text = "\n".join(texts)
    for term in REQUIRED_TERMS:
        if term not in all_text:
            problems.append(f"missing required term: {term}")
    for term in FORBIDDEN_TERMS:
        if term in all_text:
            problems.append(f"forbidden term: {term}")

    for address in sorted(set(IPV4_PATTERN.findall(all_text))):
        problems.append(f"IPv4 address exposed: {address}")
    for label, pattern in SECRET_PATTERNS:
        if pattern.search(all_text):
            problems.append(f"sensitive text exposed: {label}")

    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes, start=1):
            left = int(shape.left)
            top = int(shape.top)
            right = left + int(shape.width)
            bottom = top + int(shape.height)
            if left < 0 or top < 0 or right > slide_width or bottom > slide_height:
                problems.append(
                    f"slide {slide_index} shape {shape_index}: outside canvas "
                    f"({left}, {top}, {right}, {bottom})"
                )

    return problems


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("path", nargs="?", type=Path, default=DEFAULT_DECK)
    parser.add_argument("--dump-text", action="store_true")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    if not args.path.is_file():
        print(f"FAIL: PPTX not found: {args.path}")
        return 1

    prs = Presentation(args.path)
    if args.dump_text:
        for index, text in enumerate(slide_texts(prs), start=1):
            print(f"## Slide {index}\n{text}\n")
        return 0

    problems = validate_deck(args.path)
    if problems:
        print(f"FAIL: {len(problems)} problem(s)")
        for problem in problems:
            print(f"- {problem}")
        return 1

    print(f"PASS: {len(prs.slides)} slides")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
